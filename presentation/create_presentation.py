"""Generate a PowerPoint presentation for the Information Integration project."""

# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x1E, 0x1E, 0x2E)   # very dark blue/grey
ACCENT      = RGBColor(0x89, 0xB4, 0xFA)   # catppuccin blue
ACCENT2     = RGBColor(0xA6, 0xE3, 0xA1)   # catppuccin green
ACCENT3     = RGBColor(0xF3, 0x8B, 0xA8)   # catppuccin pink/red
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xCC, 0xC8, 0xC4)
CODE_BG     = RGBColor(0x18, 0x18, 0x2E)   # slightly darker for code boxes
SUBTITLE_C  = RGBColor(0xBA, 0xC2, 0xDE)   # muted lavender

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helper utilities ─────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_code_box(slide, code: str, left, top, width, height, font_size=10):
    """Add a dark rounded code box with monospaced text."""
    rect = add_rect(slide, left, top, width, height, CODE_BG, ACCENT, Pt(1))
    # code text
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = ACCENT2


def slide_title_bar(slide, title: str, subtitle: str = ""):
    """Add a coloured title bar at the top."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), ACCENT)
    add_text_box(slide, title,
                 Inches(0.35), Inches(0.12), Inches(11.5), Inches(0.6),
                 font_size=28, bold=True, color=DARK_BG)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(0.72), Inches(11.5), Inches(0.35),
                     font_size=14, color=DARK_BG, italic=True)


def bullet_para(tf, text: str, level: int = 0, size: int = 17,
                color: RGBColor = WHITE, bold: bool = False):
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = ("  " * level) + text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


# ── Slide builders ──────────────────────────────────────────────────────────

def make_title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, DARK_BG)

    # large accent bar
    add_rect(slide, 0, Inches(2.4), SLIDE_W, Inches(2.7), ACCENT)

    add_text_box(slide,
                 "Information Integration",
                 Inches(0.5), Inches(2.55), Inches(12.3), Inches(1.0),
                 font_size=44, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "Video Game Dataset Integration Pipeline",
                 Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6),
                 font_size=22, color=DARK_BG, align=PP_ALIGN.CENTER, italic=True)

    add_text_box(slide,
                 "Friedrich Schiller University Jena  ·  Information and Data Integration  ·  WS 2025/26",
                 Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                 font_size=13, color=SUBTITLE_C, align=PP_ALIGN.CENTER)


def make_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Project Overview", "Integrating three Kaggle video game datasets into a unified resource")

    # Three dataset cards
    card_tops = Inches(1.4)
    card_h = Inches(1.55)
    datasets = [
        ("Dataset 1", "64 017 records · 14 attributes",
         "title, console, genre, publisher,\ndeveloper, critic_score, total_sales,\nrelease_date, …"),
        ("Dataset 2", "18 800 records · 5 attributes",
         "name, platform, release_date,\nsummary, user_review"),
        ("Dataset 3", "14 055 records · 9 attributes",
         "Title, Release Date, Developer,\nPublisher, Genres, Product Rating,\nUser Score, Platforms Info"),
    ]
    card_w = Inches(3.9)
    gap = Inches(0.3)
    x = Inches(0.35)
    for ds_name, meta, attrs in datasets:
        add_rect(slide, x, card_tops, card_w, card_h, RGBColor(0x30, 0x30, 0x4E))
        add_text_box(slide, ds_name, x + Inches(0.1), card_tops + Inches(0.05),
                     card_w - Inches(0.2), Inches(0.35), font_size=15, bold=True, color=ACCENT)
        add_text_box(slide, meta, x + Inches(0.1), card_tops + Inches(0.38),
                     card_w - Inches(0.2), Inches(0.3), font_size=11, color=ACCENT3)
        add_text_box(slide, attrs, x + Inches(0.1), card_tops + Inches(0.68),
                     card_w - Inches(0.2), Inches(0.8), font_size=10, color=LIGHT_GREY)
        x += card_w + gap

    # Pipeline steps
    add_text_box(slide, "Integration Pipeline Steps", Inches(0.35), Inches(3.1),
                 Inches(12), Inches(0.35), font_size=15, bold=True, color=ACCENT)

    steps = [
        ("1", "Data Collection / Extraction", "Load CSVs via KaggleHub"),
        ("2", "Schema Mapping", "Map heterogeneous attributes to a unified target schema"),
        ("3", "Identity Resolution", "Fuzzy matching + blocking to merge duplicate records"),
        ("4", "Data Quality Assessment", "Evaluate completeness, coverage, and integration results"),
    ]
    y = Inches(3.55)
    for num, step, desc in steps:
        add_rect(slide, Inches(0.35), y, Inches(0.4), Inches(0.38), ACCENT)
        add_text_box(slide, num, Inches(0.35), y, Inches(0.4), Inches(0.38),
                     font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text_box(slide, step, Inches(0.85), y, Inches(3.2), Inches(0.38),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, desc, Inches(4.1), y, Inches(8.9), Inches(0.38),
                     font_size=12, color=LIGHT_GREY)
        y += Inches(0.52)

    # Result box
    add_rect(slide, Inches(0.35), Inches(6.05), Inches(12.6), Inches(0.38),
             RGBColor(0x1E, 0x3A, 0x1E))
    add_text_box(slide,
                 "Result: games_integrated_dataset.csv  —  unified, deduplicated, fully provenanced",
                 Inches(0.55), Inches(6.07), Inches(12.2), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT2)


def make_step1_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Step 1 – Data Collection & Extraction",
                    "Download, cache and load datasets from Kaggle using KaggleHub")

    add_text_box(slide, "Key design decisions:", Inches(0.35), Inches(1.25),
                 Inches(12), Inches(0.3), font_size=13, bold=True, color=ACCENT)

    bullets = [
        "Datasets are downloaded and cached locally via kagglehub.dataset_download()",
        "A single load_dataset() utility handles any CSV dataset in a uniform way",
        "If no target file name is given, the first CSV in the directory is loaded automatically",
        "Raw datasets are optionally persisted to data/raw/ for reproducibility",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.35), Inches(1.6), Inches(5.8), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT_GREY

    code = (
        "def load_dataset(dataset_id, target_file=None):\n"
        "    # Download & cache from Kaggle\n"
        "    path = kagglehub.dataset_download(dataset_id)\n"
        "\n"
        "    files = [f for f in os.listdir(path)\n"
        "             if f.endswith('.csv')]\n"
        "\n"
        "    if target_file:\n"
        "        if target_file not in files:\n"
        "            raise ValueError(\n"
        "                f'{target_file} not found in {files}')\n"
        "        file = target_file\n"
        "    else:\n"
        "        file = files[0]   # first CSV fallback\n"
        "\n"
        "    return pd.read_csv(os.path.join(path, file))"
    )
    add_code_box(slide, code, Inches(6.35), Inches(1.2), Inches(6.6), Inches(4.5))

    add_text_box(slide, "Pipeline call (main.py):", Inches(0.35), Inches(3.95),
                 Inches(5.8), Inches(0.3), font_size=13, bold=True, color=ACCENT)
    call_code = (
        "DATASETS = [\n"
        "    ('games_dataset_1',\n"
        "     'ujjwalaggarwal402/video-games-dataset',\n"
        "     'Video Games Data.csv'),\n"
        "    ('games_dataset_2',\n"
        "     'maso0dahmed/video-games-data', None),\n"
        "    ('games_dataset_3',\n"
        "     'beridzeg45/video-games', None),\n"
        "]\n"
        "for name, dataset_id, file in DATASETS:\n"
        "    df = load_dataset(dataset_id, file)"
    )
    add_code_box(slide, call_code, Inches(0.35), Inches(4.3), Inches(5.8), Inches(3.0))


def make_step2a_slide(prs):
    """Schema Mapping – concept & target schema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Step 2 – Schema Mapping (1/2)",
                    "Define a unified target schema and map heterogeneous attributes")

    # Left: mapping types explanation
    add_text_box(slide, "Three Mapping Types", Inches(0.35), Inches(1.25),
                 Inches(6.1), Inches(0.3), font_size=14, bold=True, color=ACCENT)

    mapping_types = [
        ("NULL mapping",
         "Drop attributes irrelevant to the target schema\n"
         "(e.g. img, na_sales, jp_sales, pal_sales, last_update)"),
        ("1:1 mapping",
         "Rename a column directly to its target name\n"
         "(e.g. 'console' → 'platform',  'name' → 'title')"),
        ("1:n mapping",
         "Expand a structured column into multiple columns\n"
         "(e.g. Platforms Info → platform + metascore per row)"),
    ]
    y = Inches(1.65)
    for mt_name, mt_desc in mapping_types:
        add_rect(slide, Inches(0.35), y, Inches(6.1), Inches(0.95),
                 RGBColor(0x28, 0x28, 0x44))
        add_text_box(slide, mt_name, Inches(0.5), y + Inches(0.05),
                     Inches(5.8), Inches(0.3), font_size=13, bold=True, color=ACCENT3)
        add_text_box(slide, mt_desc, Inches(0.5), y + Inches(0.35),
                     Inches(5.8), Inches(0.6), font_size=11, color=LIGHT_GREY)
        y += Inches(1.05)

    # entity key
    add_rect(slide, Inches(0.35), Inches(4.85), Inches(6.1), Inches(0.6),
             RGBColor(0x1E, 0x2E, 0x1E))
    add_text_box(slide,
                 "Entity key:  (title, platform)  —  same game can exist on multiple platforms",
                 Inches(0.5), Inches(4.9), Inches(5.8), Inches(0.55),
                 font_size=12, color=ACCENT2, bold=True)

    # Right: target schema table
    add_text_box(slide, "Unified Target Schema", Inches(6.8), Inches(1.25),
                 Inches(6.1), Inches(0.3), font_size=14, bold=True, color=ACCENT)

    schema_attrs = [
        "title", "platform", "release_date", "developer",
        "publisher", "genre", "critic_score", "user_score",
        "metascore", "product_rating", "total_sales",
        "summary", "source  ← provenance", "provenance  ← per-field trace",
    ]
    col_w = Inches(2.9)
    x_start = Inches(6.8)
    row_h = Inches(0.34)
    y = Inches(1.65)
    for i, attr in enumerate(schema_attrs):
        bg = RGBColor(0x2A, 0x2A, 0x44) if i % 2 == 0 else RGBColor(0x22, 0x22, 0x38)
        add_rect(slide, x_start, y, Inches(6.1), row_h, bg)
        add_text_box(slide, attr, x_start + Inches(0.1), y + Inches(0.04),
                     Inches(5.9), row_h - Inches(0.06), font_size=12, color=WHITE)
        y += row_h


def make_step2b_slide(prs):
    """Schema Mapping – apply_mapping code."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Step 2 – Schema Mapping (2/2)",
                    "apply_mapping() transforms each raw DataFrame into the target schema")

    mapping_def_code = (
        "# mappings.py – Dataset 1 correspondences\n"
        "DATASET1 = {\n"
        "    'img':          None,          # NULL\n"
        "    'title':        'title',       # 1:1\n"
        "    'console':      'platform',    # 1:1\n"
        "    'genre':        'genre',       # 1:1\n"
        "    'publisher':    'publisher',   # 1:1\n"
        "    'developer':    'developer',   # 1:1\n"
        "    'critic_score': 'critic_score',\n"
        "    'total_sales':  'total_sales',\n"
        "    'na_sales':     None,          # NULL\n"
        "    'jp_sales':     None,          # NULL\n"
        "    'pal_sales':    None,          # NULL\n"
        "    'other_sales':  None,          # NULL\n"
        "    'release_date': 'release_date',\n"
        "    'last_update':  None,          # NULL\n"
        "}"
    )
    add_code_box(slide, mapping_def_code, Inches(0.35), Inches(1.2), Inches(6.0), Inches(4.7))

    apply_code = (
        "def apply_mapping(source, df, mapping):\n"
        "    df = df.copy()\n"
        "    for orig, target in mapping.items():\n"
        "        if target is None:           # NULL\n"
        "            df.drop(columns=[orig], inplace=True)\n"
        "        elif isinstance(target, str):# 1:1\n"
        "            df.rename(columns=\n"
        "                {orig: target}, inplace=True)\n"
        "        elif isinstance(target, dict):# 1:n\n"
        "            df = extract_information(\n"
        "                df, orig, target)\n"
        "\n"
        "    df['source'] = source\n"
        "    df['provenance'] = f'origin({source})'\n"
        "    return df\n"
        "\n"
        "# 1:n: explode Platforms Info list-of-dicts\n"
        "# into separate (platform, metascore) rows"
    )
    add_code_box(slide, apply_code, Inches(6.55), Inches(1.2), Inches(6.4), Inches(4.7))

    add_text_box(slide,
                 "The 1:n mapping explodes Dataset 3's 'Platforms Info' column "
                 "(a list of {Platform, Metascore, …} dicts) into separate rows "
                 "via explode() + json_normalize(), increasing record count but preserving all platform-level data.",
                 Inches(0.35), Inches(6.0), Inches(12.6), Inches(0.5),
                 font_size=12, color=SUBTITLE_C, italic=True)


def make_step3a_slide(prs):
    """Identity Resolution – concept & preprocessing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Step 3 – Identity Resolution (1/3)",
                    "Preprocessing & normalization before matching")

    add_text_box(slide, "Why preprocessing matters", Inches(0.35), Inches(1.25),
                 Inches(12.6), Inches(0.3), font_size=14, bold=True, color=ACCENT)

    reasons = [
        ('Platform names', '"Nintendo Switch" · "Switch" · "NS" → all refer to the same platform'),
        ('Date formats',   '"2001-11-18" · "Nov 18, 2001" · "November 2001" → must be normalized to YYYY-MM-DD'),
        ('Title numbers',  '"Final Fantasy VII" → "Final Fantasy 7" (Roman numeral standardization)'),
        ('Score scales',   'Some scores are 0-10, others 0-100 → normalize to common scale'),
        ('Whitespace',     'Leading/trailing spaces cause false mismatches → strip all string fields'),
    ]
    y = Inches(1.65)
    for label, desc in reasons:
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.5),
                 RGBColor(0x28, 0x28, 0x44))
        add_text_box(slide, label, Inches(0.5), y + Inches(0.08),
                     Inches(2.1), Inches(0.35), font_size=12, bold=True, color=ACCENT3)
        add_text_box(slide, desc, Inches(2.7), y + Inches(0.08),
                     Inches(10.0), Inches(0.35), font_size=12, color=LIGHT_GREY)
        y += Inches(0.58)

    add_text_box(slide, "pre_normalize() pipeline:", Inches(0.35), Inches(4.7),
                 Inches(6.0), Inches(0.3), font_size=13, bold=True, color=ACCENT)

    pre_code = (
        "def pre_normalize(df):\n"
        "    df = normalize_dates(df)\n"
        "    df = normalize_by_mapping(df, PLATFORM, 'platform')\n"
        "    df = normalize_by_mapping(df, GENRE,    'genre')\n"
        "    df = normalize_by_mapping(df, DEV_PUB,  'developer')\n"
        "    df = normalize_by_mapping(df, DEV_PUB,  'publisher')\n"
        "    df = normalize_title_numbers(df)\n"
        "    df = normalize_scores(df)\n"
        "    df = remove_platform_all_and_missing(df)\n"
        "    df = remove_leading_trailing_whitespace(df)\n"
        "    return df"
    )
    add_code_box(slide, pre_code, Inches(0.35), Inches(5.1), Inches(6.3), Inches(2.2))

    add_text_box(slide, "LLM-assisted mapping generation:", Inches(6.85), Inches(4.7),
                 Inches(6.0), Inches(0.3), font_size=13, bold=True, color=ACCENT)
    llm_steps = (
        "1. Extract unique values for each attribute\n"
        "   (platform, genre, developer, publisher)\n"
        "   along with their frequencies\n"
        "\n"
        "2. Feed value list to ChatGPT to generate\n"
        "   normalization mapping rules\n"
        "\n"
        "3. Apply mapping rules via normalize_by_mapping()\n"
        "\n"
        "4. Manual review for edge cases & accuracy"
    )
    add_code_box(slide, llm_steps, Inches(6.85), Inches(5.1), Inches(6.1), Inches(2.2),
                 font_size=11)


def make_step3b_slide(prs):
    """Identity Resolution – matching & blocking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Step 3 – Identity Resolution (2/3)",
                    "Similarity scoring: weighted title + release date match")

    score_code = (
        "NORMALIZED_TITLE_WEIGHT = 0.85\n"
        "THRESHOLD               = 0.85\n"
        "\n"
        "def compute_match_score(row1, row2):\n"
        "    # Hard reject on platform mismatch\n"
        "    if (has_value(row1['platform']) and\n"
        "        has_value(row2['platform']) and\n"
        "        row1['platform'] != row2['platform']):\n"
        "        return 0.0\n"
        "\n"
        "    title_score = normalized_levenshtein_similarity(\n"
        "        row1['title'], row2['title'])\n"
        "\n"
        "    release_score = normalized_release_date_similarity(\n"
        "        row1['release_date'], row2['release_date'])\n"
        "\n"
        "    return (NORMALIZED_TITLE_WEIGHT * title_score +\n"
        "            (1 - NORMALIZED_TITLE_WEIGHT) * release_score)"
    )
    add_code_box(slide, score_code, Inches(0.35), Inches(1.2), Inches(6.3), Inches(4.2))

    add_text_box(slide, "Similarity components", Inches(6.65), Inches(1.25),
                 Inches(6.3), Inches(0.3), font_size=14, bold=True, color=ACCENT)
    components = [
        ("Title similarity (85%)",
         "Normalized Levenshtein distance.\n"
         "Handles small typos & abbreviations.\n"
         "e.g. 'GTA III' vs 'Grand Theft Auto 3' → low → ✗\n"
         "'Halo: Combat Evolved' vs 'Halo Combat Evolved' → high → ✓"),
        ("Release date similarity (15%)",
         "Linear decay: 1.0 if same day, 0.0 if >365 days apart.\n"
         "Missing date → 0.0 (but high title score can still\n"
         "push overall score above threshold).\n"
         "Handles regional / early-access date offsets."),
    ]
    y = Inches(1.65)
    for comp_name, comp_desc in components:
        add_rect(slide, Inches(6.65), y, Inches(6.3), Inches(1.65),
                 RGBColor(0x28, 0x28, 0x44))
        add_text_box(slide, comp_name, Inches(6.8), y + Inches(0.08),
                     Inches(6.0), Inches(0.3), font_size=13, bold=True, color=ACCENT3)
        add_text_box(slide, comp_desc, Inches(6.8), y + Inches(0.42),
                     Inches(6.0), Inches(1.2), font_size=11, color=LIGHT_GREY)
        y += Inches(1.75)

    add_text_box(slide,
                 "score(a,b)  =  0.85 · title_sim(a,b)  +  0.15 · date_sim(a,b)     ≥ 0.85  →  MATCH",
                 Inches(0.35), Inches(5.55), Inches(12.6), Inches(0.4),
                 font_size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

    date_sim_code = (
        "def normalized_release_date_similarity(r1, r2):\n"
        "    if not has_value(r1) or not has_value(r2):\n"
        "        return 0.0\n"
        "    delta = abs((r1 - r2).days)\n"
        "    return max(0.0, 1 - delta / 365)"
    )
    add_code_box(slide, date_sim_code, Inches(0.35), Inches(6.05), Inches(12.6), Inches(1.3),
                 font_size=11)


def make_step3c_slide(prs):
    """Identity Resolution – blocking & merging."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Step 3 – Identity Resolution (3/3)",
                    "Blocking + greedy matching + conflict-resolved merging")

    # Blocking diagram (left)
    add_text_box(slide, "Blocking Strategy", Inches(0.35), Inches(1.25),
                 Inches(6.1), Inches(0.3), font_size=14, bold=True, color=ACCENT)

    blocking_steps = [
        ("1  Platform block",
         "Group records by platform. Games on different\n"
         "platforms can never be the same entity."),
        ("2  Year block (±1 yr)",
         "Within each platform group, further narrow by\n"
         "release year ±1 to handle cross-year edge cases."),
        ("Edge case A",
         "Record without release date → compare against\n"
         "ALL records of the same platform block."),
        ("Edge case B",
         "Candidate without release date → added to every\n"
         "year-block of its platform so it can still match."),
        ("Greedy matching",
         "Iterate smaller dataset; each matched record\n"
         "in the larger dataset is removed from future consideration."),
    ]
    y = Inches(1.65)
    for label, desc in blocking_steps:
        add_rect(slide, Inches(0.35), y, Inches(6.1), Inches(0.85),
                 RGBColor(0x28, 0x28, 0x44))
        add_text_box(slide, label, Inches(0.5), y + Inches(0.06),
                     Inches(5.8), Inches(0.28), font_size=12, bold=True, color=ACCENT3)
        add_text_box(slide, desc, Inches(0.5), y + Inches(0.35),
                     Inches(5.8), Inches(0.48), font_size=11, color=LIGHT_GREY)
        y += Inches(0.95)

    # Merge & conflict resolution (right)
    add_text_box(slide, "Conflict Resolution Rules", Inches(6.65), Inches(1.25),
                 Inches(6.3), Inches(0.3), font_size=14, bold=True, color=ACCENT)

    rules_code = (
        "# mappings.py – resolution strategies\n"
        "ENTITY_RESOLUTION = {\n"
        "    'title':        MAX,   # longer/more complete\n"
        "    'release_date': MIN,   # earliest date\n"
        "    'developer':    MAX,\n"
        "    'publisher':    MAX,\n"
        "    'genre':        UNION, # merge genres\n"
        "    'critic_score': MAX,\n"
        "    'user_score':   MAX,\n"
        "    'summary':      UNION,\n"
        "    'total_sales':  MAX,\n"
        "    'source':       UNION, # track origins\n"
        "    'provenance':   NOTHING,\n"
        "}"
    )
    add_code_box(slide, rules_code, Inches(6.65), Inches(1.65), Inches(6.3), Inches(3.3))

    prov_code = (
        "# Provenance tracking per field\n"
        "# e.g.:\n"
        "# 'title=eq(ds1,ds2),'\n"
        "# 'genre=union(ds1,ds3),'\n"
        "# 'release_date=min(ds1,ds2)'"
    )
    add_code_box(slide, prov_code, Inches(6.65), Inches(5.1), Inches(6.3), Inches(1.5),
                 font_size=11)

    add_text_box(slide,
                 "Unmatched records from both datasets are appended as-is → outer join semantics",
                 Inches(0.35), Inches(6.6), Inches(12.6), Inches(0.35),
                 font_size=13, color=SUBTITLE_C, italic=True)


def make_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Results & Data Quality",
                    "Integration outcome and completeness assessment")

    stats = [
        ("96 872", "Total raw records\n(DS1 + DS2 + DS3)"),
        ("~68 k",  "Records after\nidentity resolution"),
        ("14",     "Attributes in\nunified schema"),
        ("100 %",  "Records with\nprovenance trace"),
    ]
    x = Inches(0.35)
    for value, label in stats:
        add_rect(slide, x, Inches(1.2), Inches(2.9), Inches(1.5),
                 RGBColor(0x28, 0x28, 0x44))
        add_text_box(slide, value, x + Inches(0.1), Inches(1.3),
                     Inches(2.7), Inches(0.65), font_size=30, bold=True,
                     color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x + Inches(0.1), Inches(1.95),
                     Inches(2.7), Inches(0.65), font_size=12,
                     color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        x += Inches(3.15)

    # Qualitative points
    add_text_box(slide, "Key observations", Inches(0.35), Inches(2.9),
                 Inches(12.6), Inches(0.3), font_size=14, bold=True, color=ACCENT)

    observations = [
        ("Deduplication",
         "Identity resolution removed duplicate records across the three datasets, "
         "producing a significantly smaller but richer unified dataset."),
        ("Schema heterogeneity",
         "Datasets used inconsistent column names, date formats, score scales, and platform labels — "
         "all resolved via explicit mapping rules and LLM-assisted normalization dictionaries."),
        ("Provenance",
         "Every field in every record tracks exactly which source dataset it came from "
         "and how conflicts were resolved (eq / min / max / union / single)."),
        ("Data Quality",
         "Completeness varies per attribute: sales data only from DS1, summaries only from DS2, "
         "product ratings only from DS3 — multi-source records are more complete."),
    ]
    y = Inches(3.3)
    for obs_title, obs_desc in observations:
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.65),
                 RGBColor(0x26, 0x26, 0x42))
        add_text_box(slide, obs_title + ":", Inches(0.5), y + Inches(0.08),
                     Inches(2.0), Inches(0.5), font_size=12, bold=True, color=ACCENT3)
        add_text_box(slide, obs_desc, Inches(2.6), y + Inches(0.08),
                     Inches(10.1), Inches(0.5), font_size=12, color=LIGHT_GREY)
        y += Inches(0.75)


def make_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    slide_title_bar(slide, "Summary", "")

    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.05), ACCENT)

    summary_items = [
        ("Data Extraction",
         "KaggleHub download + uniform CSV loading via load_dataset()"),
        ("Schema Mapping",
         "NULL / 1:1 / 1:n mappings via apply_mapping(); adds source & provenance columns"),
        ("Normalization",
         "LLM-assisted dictionaries for platforms, genres, developers; date & score standardization"),
        ("Identity Resolution",
         "Levenshtein title similarity + date proximity, with platform/year blocking"),
        ("Conflict Resolution",
         "Per-field strategies: MIN (dates), MAX (scores/text), UNION (genres), single-source fallback"),
        ("Provenance",
         "Field-level audit trail: eq(), min(), max(), union(), single() annotations"),
        ("Output",
         "Deduplicated games_integrated_dataset.csv — sorted by title, fully traceable"),
    ]

    y = Inches(1.35)
    for i, (item_title, item_desc) in enumerate(summary_items):
        bg = RGBColor(0x2C, 0x2C, 0x48) if i % 2 == 0 else RGBColor(0x24, 0x24, 0x3C)
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.67), bg)
        add_rect(slide, Inches(0.35), y, Inches(0.08), Inches(0.67), ACCENT)
        add_text_box(slide, item_title, Inches(0.55), y + Inches(0.08),
                     Inches(2.3), Inches(0.5), font_size=13, bold=True, color=ACCENT)
        add_text_box(slide, item_desc, Inches(2.95), y + Inches(0.08),
                     Inches(9.8), Inches(0.5), font_size=13, color=LIGHT_GREY)
        y += Inches(0.72)

    add_text_box(slide,
                 "Information and Data Integration  ·  FSU Jena  ·  WS 2025/26",
                 Inches(0.35), Inches(6.8), Inches(12.6), Inches(0.35),
                 font_size=11, color=SUBTITLE_C, align=PP_ALIGN.CENTER, italic=True)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_overview_slide(prs)
    make_step1_slide(prs)
    make_step2a_slide(prs)
    make_step2b_slide(prs)
    make_step3a_slide(prs)
    make_step3b_slide(prs)
    make_step3c_slide(prs)
    make_results_slide(prs)
    make_summary_slide(prs)

    output = "information_integration_presentation.pptx"
    prs.save(output)
    print(f"Saved: {output}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
