"""Generate a detailed 20-slide PowerPoint presentation for the Information Integration project."""

# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ───────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # blue
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # green
ACCENT3    = RGBColor(0xF3, 0x8B, 0xA8)   # pink/red
ACCENT4    = RGBColor(0xFA, 0xB3, 0x87)   # peach/orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xC8, 0xC4)
CODE_BG    = RGBColor(0x18, 0x18, 0x2E)
SUBTITLE_C = RGBColor(0xBA, 0xC2, 0xDE)
CARD_BG    = RGBColor(0x28, 0x28, 0x44)
ROW_A      = RGBColor(0x2A, 0x2A, 0x44)
ROW_B      = RGBColor(0x22, 0x22, 0x38)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Primitives ───────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_code(slide, code, left, top, width, height, size=10):
    add_rect(slide, left, top, width, height, CODE_BG, ACCENT, Pt(1))
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = "Courier New"
        r.font.color.rgb = ACCENT2


def title_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), ACCENT)
    add_text(slide, title, Inches(0.35), Inches(0.1),
             Inches(12.6), Inches(0.58), size=26, bold=True, color=DARK_BG)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.68),
                 Inches(12.6), Inches(0.32), size=13, italic=True, color=DARK_BG)


def section_label(slide, text, y):
    """A small coloured section heading."""
    add_text(slide, text, Inches(0.35), y, Inches(12.6), Inches(0.3),
             size=13, bold=True, color=ACCENT)


def bullet_block(slide, items, left, top, width, height,
                 size=13, color=LIGHT_GREY, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"{bullet}  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color


def info_row(slide, label, value, y, row_h=Inches(0.52),
             lw=Inches(2.8), label_color=ACCENT3, value_color=LIGHT_GREY):
    add_rect(slide, Inches(0.35), y, Inches(12.6), row_h,
             ROW_A if int(y / Inches(0.52)) % 2 == 0 else ROW_B)
    add_text(slide, label, Inches(0.5), y + Inches(0.08),
             lw - Inches(0.2), row_h - Inches(0.1),
             size=12, bold=True, color=label_color)
    add_text(slide, value, Inches(0.5) + lw, y + Inches(0.08),
             Inches(12.6) - lw - Inches(0.3), row_h - Inches(0.1),
             size=12, color=value_color)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def footer(slide, text):
    add_text(slide, text, Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.3),
             size=10, italic=True, color=SUBTITLE_C, align=PP_ALIGN.CENTER)


# ── Slide 1 – Title ──────────────────────────────────────────────────────────
def s01_title(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, Inches(2.3), SLIDE_W, Inches(2.9), ACCENT)
    add_text(slide, "Information Integration",
             Inches(0.5), Inches(2.45), Inches(12.3), Inches(1.0),
             size=46, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, "Video Game Dataset Integration Pipeline",
             Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.6),
             size=22, italic=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Friedrich Schiller University Jena  ·  Information and Data Integration  ·  WS 2025/26",
             Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
             size=13, color=SUBTITLE_C, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Objective: Integrate three heterogeneous Kaggle video game datasets into a single, "
             "deduplicated, fully provenanced unified dataset.",
             Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5),
             size=12, italic=True, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ── Slide 2 – Table of Contents ──────────────────────────────────────────────
def s02_toc(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Table of Contents", "Overview of the integration pipeline and presentation structure")

    sections = [
        ("1", "Project Motivation & Problem Statement",       "Slide 3"),
        ("2", "Data Sources",                                 "Slides 4–5"),
        ("3", "Integration Architecture",                     "Slide 6"),
        ("4", "Step 1 – Data Collection & Extraction",        "Slide 7"),
        ("5", "Step 2 – Schema Mapping",                      "Slides 8–10"),
        ("6", "Step 3 – Identity Resolution: Preprocessing",  "Slides 11–12"),
        ("7", "Step 3 – Identity Resolution: Matching",       "Slides 13–14"),
        ("8", "Step 3 – Identity Resolution: Merging",        "Slide 15"),
        ("9", "Data Quality Assessment",                      "Slides 16–17"),
        ("10", "Results & Summary",                           "Slides 18–19"),
    ]
    y = Inches(1.15)
    for num, section, pages in sections:
        bg = CARD_BG if int(num) % 2 == 0 else RGBColor(0x22, 0x22, 0x3C)
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.52), bg)
        add_rect(slide, Inches(0.35), y, Inches(0.52), Inches(0.52), ACCENT)
        add_text(slide, num, Inches(0.35), y, Inches(0.52), Inches(0.52),
                 size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(slide, section, Inches(0.97), y + Inches(0.1),
                 Inches(9.9), Inches(0.35), size=13, color=WHITE)
        add_text(slide, pages, Inches(11.0), y + Inches(0.1),
                 Inches(1.9), Inches(0.35), size=12, italic=True,
                 color=SUBTITLE_C, align=PP_ALIGN.RIGHT)
        y += Inches(0.57)


# ── Slide 3 – Motivation ─────────────────────────────────────────────────────
def s03_motivation(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Project Motivation & Problem Statement",
              "Why integrate multiple video game datasets?")

    add_text(slide, "The Challenge", Inches(0.35), Inches(1.15),
             Inches(6.1), Inches(0.3), size=14, bold=True, color=ACCENT)

    challenges = [
        "No single dataset contains all relevant video game attributes",
        "Datasets use different schemas, column names, and data formats",
        "The same game appears across datasets with inconsistent identifiers",
        "Score scales differ (0–10 vs 0–100); platform names are not standardized",
        "Release dates appear in multiple formats and may differ across regions",
        "Missing values are spread across datasets in a complementary pattern",
    ]
    bullet_block(slide, challenges, Inches(0.35), Inches(1.5),
                 Inches(6.1), Inches(3.2), size=12)

    add_text(slide, "The Goal", Inches(6.65), Inches(1.15),
             Inches(6.3), Inches(0.3), size=14, bold=True, color=ACCENT)
    goals = [
        "Produce a single unified dataset covering title, platform, release date, developer, "
        "publisher, genre, scores, sales, summary, and product rating",
        "Resolve entity identity across datasets: match records referring to the same game",
        "Merge complementary attributes: enrich each record with data from multiple sources",
        "Track provenance: record exactly where each attribute value originated",
        "Remove genuine duplicates while preserving distinct platform variants",
    ]
    bullet_block(slide, goals, Inches(6.65), Inches(1.5),
                 Inches(6.3), Inches(3.2), size=12)

    add_rect(slide, Inches(0.35), Inches(5.0), Inches(12.6), Inches(0.65),
             RGBColor(0x1E, 0x2E, 0x1E))
    add_text(slide,
             "Entity Key:  (title, platform)  —  the same game title on different platforms "
             "is treated as a distinct entity, because gameplay, release date, and scores can differ significantly.",
             Inches(0.55), Inches(5.05), Inches(12.2), Inches(0.55),
             size=12, bold=True, color=ACCENT2)

    add_text(slide, "Information Integration Methods Used", Inches(0.35), Inches(5.8),
             Inches(12.6), Inches(0.3), size=13, bold=True, color=ACCENT)
    methods = [
        "Global-as-View (GAV) schema integration — target schema defined first, sources mapped into it",
        "Fuzzy matching — Levenshtein similarity for tolerant string comparison",
        "Blocking — platform + release-year partitioning to make matching tractable",
    ]
    bullet_block(slide, methods, Inches(0.35), Inches(6.15), Inches(12.6), Inches(0.9),
                 size=12)


# ── Slide 4 – All Datasets ───────────────────────────────────────────────────
def s04_datasets(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Data Sources Overview",
              "Three Kaggle video game datasets with complementary attributes")

    datasets_info = [
        (
            "Dataset 1",
            "ujjwalaggarwal402/video-games-dataset",
            "64 017 records  ·  14 attributes",
            [
                ("img",          "NULL",          True),
                ("title",        "→ title",        False),
                ("console",      "→ platform",     False),
                ("genre",        "→ genre",        False),
                ("publisher",    "→ publisher",    False),
                ("developer",    "→ developer",    False),
                ("critic_score", "→ critic_score", False),
                ("total_sales",  "→ total_sales",  False),
                ("na_sales",     "NULL",           True),
                ("jp_sales",     "NULL",           True),
                ("pal_sales",    "NULL",           True),
                ("other_sales",  "NULL",           True),
                ("release_date", "→ release_date", False),
                ("last_update",  "NULL",           True),
            ]
        ),
        (
            "Dataset 2",
            "maso0dahmed/video-games-data",
            "18 800 records  ·  5 attributes",
            [
                ("name",         "→ title",        False),
                ("platform",     "→ platform",     False),
                ("release_date", "→ release_date", False),
                ("summary",      "→ summary",      False),
                ("user_review",  "→ user_score",   False),
            ]
        ),
        (
            "Dataset 3",
            "beridzeg45/video-games",
            "14 055 records  ·  9 attributes",
            [
                ("Title",              "→ title",          False),
                ("Release Date",       "→ release_date",   False),
                ("Developer",          "→ developer",      False),
                ("Publisher",          "→ publisher",      False),
                ("Genres",             "→ genre",          False),
                ("Product Rating",     "→ product_rating", False),
                ("User Score",         "→ user_score",     False),
                ("User Ratings Count", "NULL",             True),
                ("Platforms Info",     "1:n → platform\n+ metascore", False),
            ]
        ),
    ]

    card_w = Inches(4.08)
    gap    = Inches(0.2)
    rh     = Inches(0.33)
    card_top = Inches(1.15)

    for ci, (ds_name, kaggle_id, meta, attrs) in enumerate(datasets_info):
        x = Inches(0.35) + ci * (card_w + gap)

        # Card background
        add_rect(slide, x, card_top, card_w, Inches(6.1), CARD_BG)

        # Header block
        add_text(slide, ds_name, x + Inches(0.1), card_top + Inches(0.05),
                 card_w - Inches(0.2), Inches(0.3), size=13, bold=True, color=ACCENT)
        add_text(slide, kaggle_id, x + Inches(0.1), card_top + Inches(0.37),
                 card_w - Inches(0.2), Inches(0.25), size=8, color=SUBTITLE_C)
        add_text(slide, meta, x + Inches(0.1), card_top + Inches(0.63),
                 card_w - Inches(0.2), Inches(0.22), size=10, bold=True, color=ACCENT3)

        # Divider
        add_rect(slide, x + Inches(0.06), card_top + Inches(0.9),
                 card_w - Inches(0.12), Inches(0.02), ACCENT)

        # Attribute rows
        y = card_top + Inches(0.96)
        attr_w = card_w * 0.44
        map_w  = card_w * 0.53
        for i, (attr, mapping, is_null) in enumerate(attrs):
            bg = ROW_A if i % 2 == 0 else ROW_B
            add_rect(slide, x + Inches(0.06), y, card_w - Inches(0.12), rh, bg)
            add_text(slide, attr, x + Inches(0.1), y + Inches(0.04),
                     attr_w - Inches(0.05), rh - Inches(0.04),
                     size=10, bold=True, color=LIGHT_GREY if is_null else WHITE)
            add_text(slide, mapping, x + attr_w + Inches(0.08), y + Inches(0.04),
                     map_w - Inches(0.1), rh - Inches(0.04),
                     size=10, color=ACCENT3 if is_null else ACCENT2)
            y += rh


# ── Slide 6 – Dataset Comparison ─────────────────────────────────────────────
def s06_comparison(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Dataset Comparison & Attribute Coverage",
              "Which attributes are provided by which sources?")

    attrs = [
        "title", "platform", "release_date", "developer", "publisher",
        "genre", "critic_score", "user_score", "metascore", "product_rating",
        "total_sales", "summary",
    ]
    cols = ["Attribute", "Dataset 1", "Dataset 2", "Dataset 3"]
    coverage = {
        "title":          ("✓", "✓", "✓"),
        "platform":       ("✓", "✓", "✓ (via 1:n)"),
        "release_date":   ("✓", "✓", "✓"),
        "developer":      ("✓", "—", "✓"),
        "publisher":      ("✓", "—", "✓"),
        "genre":          ("✓", "—", "✓"),
        "critic_score":   ("✓", "—", "—"),
        "user_score":     ("—", "✓ (user_review)", "✓"),
        "metascore":      ("—", "—", "✓ (via 1:n)"),
        "product_rating": ("—", "—", "✓"),
        "total_sales":    ("✓", "—", "—"),
        "summary":        ("—", "✓", "—"),
    }

    # header
    hx = [Inches(0.35), Inches(4.0), Inches(7.0), Inches(10.2)]
    hw = [Inches(3.55), Inches(2.85), Inches(3.05), Inches(3.0)]
    hcolors = [DARK_BG, DARK_BG, DARK_BG, DARK_BG]
    add_rect(slide, Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.4), ACCENT)
    for i, col in enumerate(cols):
        add_text(slide, col, hx[i] + Inches(0.05), Inches(1.18),
                 hw[i], Inches(0.35), size=12, bold=True, color=DARK_BG)

    y = Inches(1.55)
    rh = Inches(0.42)
    for j, attr in enumerate(attrs):
        cov = coverage[attr]
        bg = ROW_A if j % 2 == 0 else ROW_B
        add_rect(slide, Inches(0.35), y, Inches(12.6), rh, bg)
        add_text(slide, attr, hx[0] + Inches(0.05), y + Inches(0.06),
                 hw[0], rh, size=12, bold=True, color=WHITE)
        for k, val in enumerate(cov):
            c = ACCENT2 if val.startswith("✓") else ACCENT3
            add_text(slide, val, hx[k+1] + Inches(0.05), y + Inches(0.06),
                     hw[k+1], rh, size=12, color=c)
        y += rh

    add_rect(slide, Inches(0.35), Inches(6.85), Inches(12.6), Inches(0.4),
             RGBColor(0x1E, 0x2E, 0x1E))
    add_text(slide,
             "Key insight: No single dataset is complete. Integration is necessary to obtain full coverage across all attributes.",
             Inches(0.55), Inches(6.88), Inches(12.2), Inches(0.35),
             size=12, bold=True, color=ACCENT2)


# ── Slide 7 – Architecture ───────────────────────────────────────────────────
def s07_architecture(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Integration Architecture",
              "Global-as-View (GAV) pipeline: extract → map → resolve → merge")

    stages = [
        ("DATA EXTRACTION",   "data_extraction.py",   "KaggleHub download\n+ CSV loading"),
        ("SCHEMA MAPPING",    "schema_mapping.py",    "NULL / 1:1 / 1:n\nattribute mapping"),
        ("NORMALIZATION",     "data_normalization.py","Dates, scores,\nplatforms, titles"),
        ("IDENTITY RES.",     "identity_resolution.py","Blocking + fuzzy\ngreedy matching"),
        ("MERGING",           "identity_resolution.py","Conflict resolution\n+ provenance"),
        ("OUTPUT",            "main.py",              "Sorted, deduplicated\nCSV dataset"),
    ]

    box_w = Inches(1.75)
    box_h = Inches(1.3)
    y_box = Inches(2.0)
    arrow_y = y_box + box_h / 2 - Inches(0.02)
    gap = Inches(0.28)   # space between boxes (includes arrow)

    x = Inches(0.4)
    for i, (stage, script, desc) in enumerate(stages):
        add_rect(slide, x, y_box, box_w, box_h, CARD_BG, ACCENT, Pt(1.5))
        add_text(slide, stage, x + Inches(0.05), y_box + Inches(0.05),
                 box_w - Inches(0.1), Inches(0.35), size=10, bold=True,
                 color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, script, x + Inches(0.05), y_box + Inches(0.4),
                 box_w - Inches(0.1), Inches(0.28), size=8.5,
                 color=ACCENT3, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.05), y_box + Inches(0.68),
                 box_w - Inches(0.1), Inches(0.55), size=9.5,
                 color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            ax = x + box_w
            add_rect(slide, ax, arrow_y, gap - Inches(0.04), Inches(0.07), ACCENT)
            add_text(slide, "▶", ax + gap - Inches(0.14), arrow_y - Inches(0.07),
                     Inches(0.18), Inches(0.2), size=9, color=ACCENT)
        x += box_w + gap

    # pairwise strategy
    add_rect(slide, Inches(0.35), Inches(3.55), Inches(12.6), Inches(1.4),
             RGBColor(0x20, 0x20, 0x3A))
    add_text(slide, "Pairwise Integration Strategy", Inches(0.55), Inches(3.6),
             Inches(12.2), Inches(0.3), size=13, bold=True, color=ACCENT)
    add_text(slide,
             "The three datasets are not merged all at once. Instead, a pairwise approach is used:\n"
             "  Step A:  DS1  ⊕  DS2  →  Intermediate Result  (merge_identities call 1)\n"
             "  Step B:  Intermediate  ⊕  DS3  →  Final Integrated Dataset  (merge_identities call 2)\n"
             "This keeps the matching logic simple and reusable: it always operates on exactly two DataFrames.",
             Inches(0.55), Inches(3.95), Inches(12.2), Inches(0.9),
             size=12, color=LIGHT_GREY)

    add_code(slide,
             "# main.py – orchestration\n"
             "df_intermediate = merge_identities(MAPPED_DATA[0][1], MAPPED_DATA[1][1])\n"
             "final_df        = merge_identities(df_intermediate,   MAPPED_DATA[2][1])\n"
             "final_df        = final_df.sort_values(by='title')\n"
             "final_df.to_csv(PROCESSED_OUTPUT_DIR / 'games_integrated_dataset.csv', index=False)",
             Inches(0.35), Inches(5.1), Inches(12.6), Inches(1.7), size=11)

    footer(slide, "GAV integration: target schema is fixed; each source is mapped into it independently")


# ── Slide 8 – Step 1: Data Extraction ────────────────────────────────────────
def s08_extraction(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 1 – Data Collection & Extraction",
              "Uniform download and loading via KaggleHub")

    section_label(slide, "Design Decisions", Inches(1.15))
    decisions = [
        "All three datasets are downloaded and cached locally via kagglehub.dataset_download()",
        "A single generic load_dataset() function handles any Kaggle CSV dataset uniformly",
        "If no target filename is supplied, the first CSV found in the directory is used as a fallback",
        "Raw datasets are saved to data/raw/ for reproducibility and manual inspection",
        "The pipeline decouples download from processing: load once, process many times",
    ]
    bullet_block(slide, decisions, Inches(0.35), Inches(1.48),
                 Inches(6.1), Inches(2.3), size=12)

    add_code(slide,
             "def load_dataset(dataset_id, target_file=None):\n"
             "    \"\"\"Download, cache and load a Kaggle CSV dataset.\"\"\"\n"
             "    path = kagglehub.dataset_download(dataset_id)\n"
             "\n"
             "    files = [f for f in os.listdir(path)\n"
             "             if f.endswith('.csv')]\n"
             "\n"
             "    if target_file:\n"
             "        if target_file not in files:\n"
             "            raise ValueError(\n"
             "                f'{target_file} not found in {files}')\n"
             "        file = target_file\n"
             "    else:\n"
             "        file = files[0]   # first CSV as fallback\n"
             "\n"
             "    return pd.read_csv(os.path.join(path, file))",
             Inches(6.35), Inches(1.15), Inches(6.6), Inches(4.5))

    section_label(slide, "Pipeline invocation (main.py)", Inches(3.9))
    add_code(slide,
             "DATASETS = [\n"
             "    ('games_dataset_1',\n"
             "     'ujjwalaggarwal402/video-games-dataset',\n"
             "     'Video Games Data.csv'),   # explicit filename\n"
             "    ('games_dataset_2',\n"
             "     'maso0dahmed/video-games-data', None),   # auto\n"
             "    ('games_dataset_3',\n"
             "     'beridzeg45/video-games', None),          # auto\n"
             "]\n"
             "\n"
             "for name, dataset_id, file in DATASETS:\n"
             "    df = load_dataset(dataset_id, file)\n"
             "    df.to_csv(RAW_OUTPUT_DIR / f'{name}.csv', index=False)",
             Inches(0.35), Inches(4.28), Inches(6.1), Inches(3.05))


# ── Slide 9 – Step 2: Schema Mapping Concepts ────────────────────────────────
def s09_schema_concepts(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 2 – Schema Mapping: Concepts",
              "Global-as-View: define the target schema first, then map each source into it")

    add_text(slide, "Integration Approach: Global-as-View (GAV)",
             Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.3),
             size=14, bold=True, color=ACCENT)
    add_text(slide,
             "The unified target schema is defined upfront. Each source dataset is then mapped into this schema "
             "independently. This approach cleanly separates the schema definition from the mapping logic "
             "and makes it easy to add new sources in the future.",
             Inches(0.35), Inches(1.5), Inches(12.6), Inches(0.65),
             size=12, color=LIGHT_GREY)

    mapping_types = [
        ("NULL mapping",
         "Drop attributes that are not relevant to the target schema.\n"
         "Examples: img (image URL), na_sales / jp_sales / pal_sales (regional breakdowns), last_update.\n"
         "Implementation: df.drop(columns=[col], inplace=True)",
         ACCENT3),
        ("1:1 mapping",
         "Directly rename a source column to its target name without any transformation.\n"
         "Examples: 'console' → 'platform',  'name' → 'title',  'Release Date' → 'release_date'.\n"
         "Implementation: df.rename(columns={orig: target}, inplace=True)",
         ACCENT),
        ("1:n mapping",
         "Expand one structured source column into multiple target columns (and potentially more rows).\n"
         "Example: Dataset 3's 'Platforms Info' is a JSON-like list of dicts "
         "(one per platform + metascore).\n"
         "→ df.explode() + pd.json_normalize() creates one row per (title, platform) combination.",
         ACCENT2),
    ]

    y = Inches(2.3)
    for mt_name, mt_desc, color in mapping_types:
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(1.3),
                 RGBColor(0x26, 0x26, 0x42))
        add_rect(slide, Inches(0.35), y, Inches(0.1), Inches(1.3), color)
        add_text(slide, mt_name, Inches(0.55), y + Inches(0.08),
                 Inches(12.2), Inches(0.3), size=13, bold=True, color=color)
        add_text(slide, mt_desc, Inches(0.55), y + Inches(0.42),
                 Inches(12.2), Inches(0.8), size=11, color=LIGHT_GREY)
        y += Inches(1.4)

    add_rect(slide, Inches(0.35), Inches(6.55), Inches(12.6), Inches(0.45),
             RGBColor(0x1E, 0x2A, 0x1E))
    add_text(slide,
             "Source & provenance columns are added automatically by apply_mapping() to every record: "
             "source = dataset name,  provenance = 'origin(<source>)'  (updated during merging)",
             Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.35),
             size=12, bold=True, color=ACCENT2)


# ── Slide 10 – Step 2: apply_mapping code ────────────────────────────────────
def s10_apply_mapping(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 2 – Schema Mapping: apply_mapping()",
              "Core function that transforms any source DataFrame into the target schema")

    add_code(slide,
             "def apply_mapping(source, df, mapping):\n"
             "    df = df.copy()\n"
             "\n"
             "    for orig_attr, target_attr in mapping.items():\n"
             "        if target_attr is None:          # NULL\n"
             "            df.drop(columns=[orig_attr],\n"
             "                    inplace=True)\n"
             "\n"
             "        elif isinstance(target_attr, str):  # 1:1\n"
             "            df.rename(\n"
             "                columns={orig_attr: target_attr},\n"
             "                inplace=True)\n"
             "\n"
             "        elif isinstance(target_attr, dict): # 1:n\n"
             "            df = extract_information(\n"
             "                df, orig_attr, target_attr)\n"
             "\n"
             "        else:\n"
             "            raise ValueError(\n"
             "                f\"Invalid mapping: {orig_attr}\")\n"
             "\n"
             "    # Provenance metadata\n"
             "    df['source']     = source\n"
             "    df['provenance'] = f'origin({source})'\n"
             "    return df",
             Inches(0.35), Inches(1.15), Inches(6.3), Inches(5.5))

    add_code(slide,
             "def extract_information(df, orig_attr, target_attr):\n"
             "    \"\"\"1:n mapping: explode list-of-dicts column.\"\"\"\n"
             "    df = df.copy()\n"
             "\n"
             "    # Parse string repr of list-of-dicts\n"
             "    df[orig_attr] = df[orig_attr].apply(\n"
             "        lambda x:\n"
             "            ast.literal_eval(x) if pd.notna(x) else [])\n"
             "\n"
             "    # One row per list element\n"
             "    df = df.explode(orig_attr)\n"
             "\n"
             "    # Normalise dicts into columns\n"
             "    extracted = df[orig_attr].apply(\n"
             "        lambda x: x if isinstance(x, dict) else {})\n"
             "    extracted_df = pd.json_normalize(extracted)\n"
             "\n"
             "    # Rename to target names, drop unmapped\n"
             "    rename_map = {k: v\n"
             "        for k, v in target_attr.items() if v}\n"
             "    extracted_df = (extracted_df\n"
             "        .rename(columns=rename_map)\n"
             "        [list(rename_map.values())])\n"
             "\n"
             "    df = df.drop(columns=[orig_attr])\n"
             "           .reset_index(drop=True)\n"
             "    return pd.concat(\n"
             "        [df, extracted_df.reset_index(drop=True)],\n"
             "        axis=1)",
             Inches(6.75), Inches(1.15), Inches(6.2), Inches(5.5))

    footer(slide,
           "Dataset 3: 14 055 source rows  →  ~52 000 rows after 1:n expansion (one row per title × platform)")


# ── Slide 11 – Step 2: Mapping Definitions ───────────────────────────────────
def s11_mapping_defs(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 2 – Schema Mapping: Mapping Definitions",
              "Explicit correspondence dictionaries in mappings.py")

    add_code(slide,
             "# Dataset 1 mapping\n"
             "DATASET1 = {\n"
             "    'img':          None,\n"
             "    'title':        'title',\n"
             "    'console':      'platform',\n"
             "    'genre':        'genre',\n"
             "    'publisher':    'publisher',\n"
             "    'developer':    'developer',\n"
             "    'critic_score': 'critic_score',\n"
             "    'total_sales':  'total_sales',\n"
             "    'na_sales':     None,\n"
             "    'jp_sales':     None,\n"
             "    'pal_sales':    None,\n"
             "    'other_sales':  None,\n"
             "    'release_date': 'release_date',\n"
             "    'last_update':  None,\n"
             "}",
             Inches(0.35), Inches(1.15), Inches(4.0), Inches(5.2))

    add_code(slide,
             "# Dataset 2 mapping\n"
             "DATASET2 = {\n"
             "    'name':         'title',\n"
             "    'platform':     'platform',\n"
             "    'release_date': 'release_date',\n"
             "    'summary':      'summary',\n"
             "    'user_review':  'user_score',\n"
             "}",
             Inches(4.5), Inches(1.15), Inches(4.0), Inches(2.6))

    add_code(slide,
             "# Dataset 3 mapping\n"
             "DATASET3 = {\n"
             "    'Title':        'title',\n"
             "    'Release Date': 'release_date',\n"
             "    'Developer':    'developer',\n"
             "    'Publisher':    'publisher',\n"
             "    'Genres':       'genre',\n"
             "    'Product Rating':'product_rating',\n"
             "    'User Score':   'user_score',\n"
             "    'User Ratings Count': None,\n"
             "    'Platforms Info': {\n"
             "        'Platform':          'platform',\n"
             "        'Platform Metascore':'metascore',\n"
             "        'Platform Metascore Count': None,\n"
             "    },\n"
             "}",
             Inches(4.5), Inches(3.9), Inches(4.0), Inches(3.75))

    add_code(slide,
             "# Main pipeline\n"
             "MAPPINGS = {\n"
             "    'games_dataset_1': mappings.DATASET1,\n"
             "    'games_dataset_2': mappings.DATASET2,\n"
             "    'games_dataset_3': mappings.DATASET3,\n"
             "}\n"
             "\n"
             "MAPPED_DATA = []\n"
             "for name, df in RAW_DATA:\n"
             "    mapped_df = apply_mapping(\n"
             "        name, df, MAPPINGS[name])\n"
             "    MAPPED_DATA.append((name, mapped_df))",
             Inches(8.7), Inches(1.15), Inches(4.28), Inches(3.5))

    add_rect(slide, Inches(8.7), Inches(4.8), Inches(4.28), Inches(1.55),
             RGBColor(0x20, 0x20, 0x3A))
    add_text(slide, "Why explicit mapping dicts?", Inches(8.85), Inches(4.88),
             Inches(4.0), Inches(0.3), size=12, bold=True, color=ACCENT)
    notes = [
        "Declarative — easy to audit and extend",
        "Decoupled from logic in apply_mapping()",
        "Version-controllable as plain Python",
        "Serves as documentation of schema decisions",
    ]
    bullet_block(slide, notes, Inches(8.85), Inches(5.2), Inches(4.0), Inches(1.1),
                 size=11)


# ── Slide 12 – Step 3: Preprocessing Overview ────────────────────────────────
def s12_preprocessing(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 3 – Identity Resolution: Preprocessing",
              "Normalize data before matching to reduce false negatives")

    steps = [
        ("normalize_dates()",
         "Converts all release_date values to datetime.date (YYYY-MM-DD).\n"
         "Handles: various string formats, 'tbd' → pd.NA, missing → pd.NA, timezone-aware parsing.",
         "pd.to_datetime(df['release_date'], errors='coerce', utc=True).dt.date"),
        ("normalize_by_mapping()",
         "Maps platform, genre, developer, publisher values to canonical forms using predefined dicts.\n"
         "Example: 'winp' → 'PC',  'ps4' → 'PlayStation 4',  'nin switch' → 'Nintendo Switch'.",
         "df[col].str.lower().map(mapping).fillna(df[col])"),
        ("normalize_title_numbers()",
         "Converts Roman numerals in titles to Arabic: 'Final Fantasy VII' → 'Final Fantasy 7'.\n"
         "Uses the roman library; handles edge cases like 'IIII' → '4'.",
         "ROMAN_PATTERN.sub(lambda m: str(roman.fromRoman(m.group(0))), title)"),
        ("normalize_scores()",
         "Replaces 'tbd' (case-insensitive) with pd.NA in critic_score, user_score, metascore columns.\n"
         "Score scale differences across datasets are tolerated at matching time.",
         "df[col].replace(r'(?i)^\\s*tbd\\s*$', pd.NA, regex=True)"),
        ("remove_platform_all_and_missing()",
         "Drops records where platform is 'All' or missing, as they cannot be matched by platform.\n"
         "Such records introduce noise without contributing to identity resolution.",
         "df[df['platform'].notna() & (df['platform'] != 'All')]"),
        ("remove_leading_trailing_whitespace()",
         "Strips whitespace from all string columns to avoid trivial mismatches.",
         "df[col].str.strip()  for all object columns"),
    ]

    y = Inches(1.15)
    rh = Inches(0.97)
    for fn_name, desc, code_snippet in steps:
        add_rect(slide, Inches(0.35), y, Inches(12.6), rh, CARD_BG)
        add_text(slide, fn_name, Inches(0.5), y + Inches(0.06),
                 Inches(2.9), Inches(0.3), size=12, bold=True, color=ACCENT3)
        add_text(slide, desc, Inches(3.5), y + Inches(0.06),
                 Inches(5.6), Inches(0.52), size=10, color=LIGHT_GREY)
        add_rect(slide, Inches(9.2), y + Inches(0.08), Inches(3.6), rh - Inches(0.16),
                 CODE_BG, ACCENT, Pt(0.5))
        add_text(slide, code_snippet, Inches(9.3), y + Inches(0.14),
                 Inches(3.4), rh - Inches(0.2), size=8.5,
                 color=ACCENT2)
        y += rh + Inches(0.04)


# ── Slide 13 – LLM-Assisted Normalization ────────────────────────────────────
def s13_llm_mapping(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 3 – LLM-Assisted Normalization Mapping",
              "How normalization dictionaries were generated for platform, genre, developer, publisher")

    add_text(slide, "The Challenge", Inches(0.35), Inches(1.15),
             Inches(12.6), Inches(0.3), size=14, bold=True, color=ACCENT)
    add_text(slide,
             "Platform names across three datasets are inconsistent: 'Nintendo Switch', 'Switch', 'NS', 'nintendo switch', "
             "'NSwitch' all refer to the same platform. Manually curating thousands of unique values is error-prone and "
             "time-consuming. The same problem applies to genre names, developer/publisher names.",
             Inches(0.35), Inches(1.5), Inches(12.6), Inches(0.7),
             size=12, color=LIGHT_GREY)

    add_text(slide, "LLM-Assisted Workflow", Inches(0.35), Inches(2.3),
             Inches(6.1), Inches(0.3), size=13, bold=True, color=ACCENT)

    steps = [
        ("1. Extract unique values",
         "Use pandas to collect all unique values for each attribute "
         "(platform, genre, developer, publisher) across all three datasets, "
         "together with their occurrence frequencies."),
        ("2. Prompt ChatGPT",
         "Feed the sorted frequency list to ChatGPT (GPT-4o) and ask it to generate "
         "a normalization mapping: raw value → canonical name. The frequency context "
         "helps the model prioritise the most common canonical form."),
        ("3. Apply the mapping",
         "Store the output as a Python dict in mappings.py. "
         "normalize_by_mapping() applies it: str.lower().map(mapping).fillna(original)."),
        ("4. Manual review",
         "Inspect edge cases flagged by mismatches. Correct ambiguous entries manually "
         "(e.g. 'I' as Roman numeral vs common word, abbreviations shared across platforms)."),
    ]
    y = Inches(2.65)
    for step_title, step_desc in steps:
        add_rect(slide, Inches(0.35), y, Inches(6.1), Inches(0.95), CARD_BG)
        add_text(slide, step_title, Inches(0.5), y + Inches(0.07),
                 Inches(5.8), Inches(0.28), size=12, bold=True, color=ACCENT3)
        add_text(slide, step_desc, Inches(0.5), y + Inches(0.38),
                 Inches(5.8), Inches(0.52), size=11, color=LIGHT_GREY)
        y += Inches(1.02)

    add_code(slide,
             "# Excerpt from PLATFORM mapping (mappings.py)\n"
             "PLATFORM = {\n"
             "    'pc':              'PC',\n"
             "    'winp':            'PC',\n"
             "    'windows':         'PC',\n"
             "    'playstation':     'PlayStation',\n"
             "    'ps':              'PlayStation',\n"
             "    'ps1':             'PlayStation',\n"
             "    'nintendo switch': 'Nintendo Switch',\n"
             "    'switch':          'Nintendo Switch',\n"
             "    'ns':              'Nintendo Switch',\n"
             "    'osx':             'MacOS',\n"
             "    'mac':             'MacOS',\n"
             "    ...\n"
             "}",
             Inches(6.65), Inches(2.3), Inches(6.3), Inches(3.5))

    add_rect(slide, Inches(6.65), Inches(5.95), Inches(6.3), Inches(0.65),
             RGBColor(0x20, 0x20, 0x3A))
    add_text(slide,
             "Benefit: a large fraction of mismatches are resolved automatically before any "
             "record-level comparison takes place, significantly improving match recall.",
             Inches(6.8), Inches(6.0), Inches(6.0), Inches(0.55),
             size=11, italic=True, color=SUBTITLE_C)


# ── Slide 14 – Matching Logic ─────────────────────────────────────────────────
def s14_matching(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 3 – Identity Resolution: Matching Logic",
              "Weighted similarity score: title (85%) + release date (15%)")

    add_code(slide,
             "TITLE_WEIGHT = 0.85\n"
             "THRESHOLD    = 0.85\n"
             "\n"
             "def compute_match_score(row1, row2):\n"
             "    # Hard reject: platform mismatch\n"
             "    if (has_value(row1['platform']) and\n"
             "        has_value(row2['platform']) and\n"
             "        row1['platform'] != row2['platform']):\n"
             "        return 0.0\n"
             "\n"
             "    title_score = normalized_levenshtein_similarity(\n"
             "        row1.get('title'), row2.get('title'))\n"
             "\n"
             "    release_score = normalized_release_date_similarity(\n"
             "        row1.get('release_date'), row2.get('release_date'))\n"
             "\n"
             "    return TITLE_WEIGHT * title_score + (1 - TITLE_WEIGHT) * release_score\n"
             "\n"
             "def normalized_levenshtein_similarity(a, b):\n"
             "    if pd.isna(a) or pd.isna(b): return 0.0\n"
             "    return ratio(str(a), str(b))   # Levenshtein.ratio\n"
             "\n"
             "def normalized_release_date_similarity(r1, r2):\n"
             "    if not has_value(r1) or not has_value(r2): return 0.0\n"
             "    if not isinstance(r1, date) or not isinstance(r2, date): return 0.0\n"
             "    delta_days = abs((r1 - r2).days)\n"
             "    return max(0.0, 1 - delta_days / 365)",
             Inches(0.35), Inches(1.15), Inches(6.5), Inches(5.3))

    add_text(slide, "Why this design?", Inches(7.0), Inches(1.15),
             Inches(6.0), Inches(0.3), size=14, bold=True, color=ACCENT)

    reasons = [
        ("85% title weight",
         "Title is the primary identifier. High weight ensures only "
         "genuinely similar titles match."),
        ("15% date weight",
         "Breaks ties between similar titles (e.g. 'GTA 3' vs 'GTA 4') "
         "that are clearly different games with different release years."),
        ("Levenshtein distance",
         "Handles minor differences: punctuation, spacing, "
         "abbreviations, subtitle variations."),
        ("Linear date decay",
         "1.0 if same day → 0.0 if >365 days apart. "
         "Tolerates regional / early-access date offsets."),
        ("Missing date → 0.0",
         "Conservative: no date evidence contributes 0. "
         "High title score can still exceed threshold alone (0.85 × 1.0 = 0.85 ✓)."),
        ("Platform hard reject",
         "Same game on different platforms is a different entity "
         "by definition → instant reject without score calculation."),
        ("Threshold = 0.85",
         "Empirically chosen to balance precision and recall. "
         "High enough to avoid false positives."),
    ]
    y = Inches(1.5)
    for r_title, r_desc in reasons:
        add_rect(slide, Inches(7.0), y, Inches(6.0), Inches(0.72), CARD_BG)
        add_text(slide, r_title, Inches(7.15), y + Inches(0.06),
                 Inches(5.7), Inches(0.25), size=11, bold=True, color=ACCENT3)
        add_text(slide, r_desc, Inches(7.15), y + Inches(0.3),
                 Inches(5.7), Inches(0.38), size=11, color=LIGHT_GREY)
        y += Inches(0.79)

    add_text(slide,
             "score(a,b)  =  0.85 · lev_sim(title_a, title_b)  +  0.15 · date_sim(date_a, date_b)     ≥ 0.85  →  MATCH",
             Inches(0.35), Inches(7), Inches(12.6), Inches(0.4),
             size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


# ── Slide 15 – Blocking Strategy ─────────────────────────────────────────────
def s15_blocking(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 3 – Identity Resolution: Blocking Strategy",
              "Reduce O(n²) comparisons to a manageable number")

    add_text(slide, "Why blocking is necessary", Inches(0.35), Inches(1.15),
             Inches(12.6), Inches(0.3), size=13, bold=True, color=ACCENT)
    add_text(slide,
             "Naively comparing every record in DS1 against every record in DS2 would require "
             "64 017 × 18 800 ≈ 1.2 billion comparisons — computationally infeasible. "
             "Blocking partitions the data so each record is only compared against a small, "
             "relevant subset.",
             Inches(0.35), Inches(1.5), Inches(6.6), Inches(0.6),
             size=12, color=LIGHT_GREY)

    add_text(slide, "Two-Level Blocking", Inches(0.35), Inches(2.2),
             Inches(6.1), Inches(0.3), size=13, bold=True, color=ACCENT)

    blocks = [
        ("Level 1: Platform block",
         "Group all records by their normalized platform name. "
         "Games on different platforms cannot be the same entity → no comparison needed.",
         ACCENT),
        ("Level 2: Year block (±1 yr)",
         "Within each platform group, further partition by release year. "
         "Include years year−1, year, year+1 in one block to handle edge cases "
         "(e.g. 2002-12-30 vs 2003-01-02).",
         ACCENT2),
        ("Edge case A – no date (query side)",
         "If the record being matched has no release date, fall back to the "
         "full platform-only block (all records of that platform).",
         ACCENT3),
        ("Edge case B – no date (candidate side)",
         "If a candidate in the larger dataset has no date, add it to every "
         "year-block of its platform so it can still be matched against any record.",
         ACCENT4),
    ]
    y = Inches(2.55)
    for b_title, b_desc, color in blocks:
        add_rect(slide, Inches(0.35), y, Inches(6.1), Inches(0.98), CARD_BG)
        add_rect(slide, Inches(0.35), y, Inches(0.12), Inches(0.98), color)
        add_text(slide, b_title, Inches(0.55), y + Inches(0.07),
                 Inches(5.7), Inches(0.28), size=12, bold=True, color=color)
        add_text(slide, b_desc, Inches(0.55), y + Inches(0.38),
                 Inches(5.7), Inches(0.55), size=11, color=LIGHT_GREY)
        y += Inches(1.06)

    add_code(slide,
             "# Build platform+year blocks from the larger dataset\n"
             "for idx, row in enumerate(df2_records):\n"
             "    platform = row.get('platform')\n"
             "    date_val  = row.get('release_date', pd.NA)\n"
             "\n"
             "    platform_only_blocks.setdefault(platform, {'indices': []})\\\n"
             "        ['indices'].append(idx)\n"
             "\n"
             "    if has_value(date_val) and isinstance(date_val, date):\n"
             "        for y in (date_val.year-1, date_val.year, date_val.year+1):\n"
             "            platform_year_blocks.setdefault(\n"
             "                (platform, y), {'indices': []})\\\n"
             "                ['indices'].append(idx)\n"
             "    else:   # edge case B\n"
             "        missing_date_by_platform\\\n"
             "            .setdefault(platform, []).append(idx)\n"
             "\n"
             "# Add missing-date rows to every year block of their platform\n"
             "for plat, idxs in missing_date_by_platform.items():\n"
             "    for key, block in platform_year_blocks.items():\n"
             "        if key[0] == plat:\n"
             "            block['indices'].extend(idxs)",
             Inches(6.88), Inches(1.24), Inches(6.3), Inches(5.5))

    add_text(slide, "Greedy matching: each matched record in the larger dataset is "
             "removed from consideration → prevents duplicate assignments",
             Inches(0.35), Inches(6.75), Inches(6.0), Inches(0.45),
             size=11, italic=True, color=SUBTITLE_C)


# ── Slide 16 – Merging & Conflict Resolution ─────────────────────────────────
def s16_merging(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Step 3 – Identity Resolution: Merging & Conflict Resolution",
              "How matched records are combined into a single enriched record")

    add_code(slide,
             "# Resolution strategies (mappings.py)\n"
             "ENTITY_RESOLUTION = {\n"
             "    'title':         MAX,    # longer = more descriptive\n"
             "    'release_date':  MIN,    # earliest = most likely correct\n"
             "    'developer':     MAX,\n"
             "    'publisher':     MAX,\n"
             "    'genre':         UNION,  # combine genre lists\n"
             "    'platform':      NOTHING,# assumed equal by matching\n"
             "    'critic_score':  MAX,\n"
             "    'user_score':    MAX,\n"
             "    'metascore':     MAX,\n"
             "    'summary':       UNION,\n"
             "    'product_rating':MAX,\n"
             "    'total_sales':   MAX,\n"
             "    'source':        UNION,  # list all contributing sources\n"
             "    'provenance':    NOTHING,# handled separately\n"
             "}",
             Inches(0.35), Inches(1.15), Inches(6.1), Inches(3.2))

    add_code(slide,
             "def merge_records(row1, row2):\n"
             "    merged = {'provenance': ''}\n"
             "    src1, src2 = row1['source'], row2['source']\n"
             "    all_cols = (set(row1) | set(row2)) - {'provenance'}\n"
             "    for col in all_cols:\n"
             "        v1 = row1.get(col, pd.NA)\n"
             "        v2 = row2.get(col, pd.NA)\n"
             "        if has_value(v1) and has_value(v2):\n"
             "            if v1 == v2:\n"
             "                merged[col] = v1\n"
             "                prov = f'{col}=eq({src1},{src2})'\n"
             "            else:\n"
             "                res = ENTITY_RESOLUTION.get(col, NOTHING)\n"
             "                if res == MIN:\n"
             "                    merged[col] = min_columns(v1, v2)\n"
             "                    prov = f'{col}=min({src1},{src2})'\n"
             "                elif res == MAX:\n"
             "                    merged[col] = max_columns(v1, v2)\n"
             "                    prov = f'{col}=max({src1},{src2})'\n"
             "                elif res == UNION:\n"
             "                    merged[col] = union_delimited(v1, v2)\n"
             "                    prov = f'{col}=union({src1},{src2})'\n"
             "                else:  # NOTHING\n"
             "                    merged[col] = v1; prov = ''\n"
             "            if prov:\n"
             "                merged['provenance'] = union_delimited(\n"
             "                    merged['provenance'], prov, ',')\n"
             "        elif has_value(v1):\n"
             "            merged[col] = v1\n"
             "            merged['provenance'] += f',{col}=single({src1})'\n"
             "        elif has_value(v2):\n"
             "            merged[col] = v2\n"
             "            merged['provenance'] += f',{col}=single({src2})'\n"
             "        else:\n"
             "            merged[col] = pd.NA\n"
             "    return merged",
             Inches(6.55), Inches(1.15), Inches(6.43), Inches(5.5), size=8.5)

    add_rect(slide, Inches(0.35), Inches(4.5), Inches(5.85), Inches(1.45),
             RGBColor(0x20, 0x20, 0x3A))
    add_text(slide, "Provenance example (output field):", Inches(0.5), Inches(4.58),
             Inches(5.5), Inches(0.28), size=12, bold=True, color=ACCENT)
    add_text(slide,
             "title=max(ds1,ds2), release_date=min(ds1,ds2),\n"
             "genre=union(ds1,ds3), summary=single(ds2),\n"
             "critic_score=single(ds1), user_score=eq(ds2,ds3)",
             Inches(0.5), Inches(4.88), Inches(5.5), Inches(0.95),
             size=11, color=ACCENT2)


# ── Slide 17 – Data Quality: Overview ────────────────────────────────────────
def s17_quality_overview(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Data Quality Assessment (1/2)",
              "Integration results and identified data quality issues")

    stats = [
        ("96 872", "Raw records\n(DS1+DS2+DS3 combined)"),
        ("80 523", "Final integrated\ndataset size"),
        ("16 349", "Records merged\n(cross-source pairs)"),
        ("100%",   "Records with\nprovenance metadata"),
    ]
    x = Inches(0.35)
    for val, label in stats:
        add_rect(slide, x, Inches(1.15), Inches(3.0), Inches(1.15), CARD_BG)
        add_text(slide, val, x + Inches(0.1), Inches(1.2),
                 Inches(2.8), Inches(0.55), size=30, bold=True,
                 color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Inches(0.1), Inches(1.75),
                 Inches(2.8), Inches(0.48), size=11, color=LIGHT_GREY,
                 align=PP_ALIGN.CENTER)
        x += Inches(3.25)

    section_label(slide, "Missing & Sparse Attributes", Inches(2.45))
    issues = [
        ("critic_score",
         "Only in Dataset 1 → missing for all records that only exist in DS2 or DS3"),
        ("summary",
         "Only in Dataset 2 → missing for unmatched DS1 and DS3 records"),
        ("total_sales",
         "Only in Dataset 1 → missing for DS2/DS3-only records"),
        ("product_rating",
         "Only in Dataset 3 → missing for DS1/DS2-only records"),
        ("platform-level summary",
         "Same game on different platforms may have different or missing summaries; "
         "these cannot be propagated across platforms (entity key = title + platform)"),
    ]
    y = Inches(2.8)
    for attr, desc in issues:
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.55),
                 ROW_A if issues.index((attr, desc)) % 2 == 0 else ROW_B)
        add_text(slide, attr, Inches(0.5), y + Inches(0.1),
                 Inches(2.3), Inches(0.35), size=12, bold=True, color=ACCENT3)
        add_text(slide, desc, Inches(2.9), y + Inches(0.1),
                 Inches(9.8), Inches(0.35), size=12, color=LIGHT_GREY)
        y += Inches(0.58)

    add_text(slide, "Provenance in action:", Inches(0.35), Inches(5.82),
             Inches(12.6), Inches(0.28), size=13, bold=True, color=ACCENT)
    add_text(slide,
             "When a data quality issue is found in the integrated dataset, the provenance field "
             "allows reverse-engineering: which source provided this value, and which resolution "
             "strategy was applied? This helps distinguish errors in source data from integration artifacts.",
             Inches(0.35), Inches(6.13), Inches(12.6), Inches(0.55),
             size=12, color=LIGHT_GREY)


# ── Slide 18 – Data Quality: Known Issues ────────────────────────────────────
def s18_quality_issues(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Data Quality Assessment (2/2)",
              "Known data quality issues and their root causes")

    issues = [
        ("Duplicates from within-dataset repetitions",
         "Dataset 2 contains multiple entries for the same (title, platform), e.g. '4x4 Evo 2, PC' "
         "appears twice (EU and NA release). Since identity resolution operates across datasets "
         "(not within), these survive as duplicates in the final output.",
         ACCENT3),
        ("1:n expansion date propagation",
         "Dataset 3's Platforms Info lists multiple platforms but assigns a single release date to all. "
         "After explode(), all platform variants share the same date even if actual regional dates differ. "
         "This can cause false positives in date-based matching.",
         ACCENT4),
        ("Genre heterogeneity across platforms",
         "The same game may be tagged with different genres on different platforms across datasets. "
         "Since entity key = (title, platform), these cannot be merged → the integrated dataset "
         "may show different genre values for the same logical game on different platforms.",
         ACCENT),
        ("Multi-year early-access gaps",
         "Games with long early-access periods (e.g. 'Mount & Blade 2: Bannerlord', 2-year gap) "
         "may appear with two very different release dates. The ±1-year blocking may not cover these, "
         "leading to unmatched duplicates.",
         ACCENT2),
        ("Roman numeral edge cases",
         "Common words that happen to be valid Roman numerals (e.g., 'I', 'V', 'C') are candidates "
         "for incorrect conversion. The implementation uses the roman library which rejects most "
         "non-numeral tokens, but edge cases like 'IIII' require special-case handling.",
         SUBTITLE_C),
    ]
    y = Inches(1.15)
    for i_title, i_desc, color in issues:
        add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(1.08),
                 CARD_BG)
        add_rect(slide, Inches(0.35), y, Inches(0.1), Inches(1.08), color)
        add_text(slide, i_title, Inches(0.55), y + Inches(0.07),
                 Inches(12.2), Inches(0.28), size=12, bold=True, color=color)
        add_text(slide, i_desc, Inches(0.55), y + Inches(0.38),
                 Inches(12.2), Inches(0.65), size=11, color=LIGHT_GREY)
        y += Inches(1.14)


# ── Slide 19 – Results ───────────────────────────────────────────────────────
def s19_results(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Results",
              "Final integrated dataset: games_integrated_dataset.csv")

    # Stats row
    stats = [
        ("80 523",  "Unified records"),
        ("14",      "Schema attributes"),
        ("2-step",  "Pairwise merges"),
        ("~17%",    "Records enriched\nfrom 2+ sources"),
    ]
    x = Inches(0.35)
    for val, label in stats:
        add_rect(slide, x, Inches(1.15), Inches(3.0), Inches(1.0), CARD_BG)
        add_text(slide, val, x + Inches(0.1), Inches(1.18),
                 Inches(2.8), Inches(0.5), size=26, bold=True,
                 color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Inches(0.1), Inches(1.68),
                 Inches(2.8), Inches(0.4), size=11, color=LIGHT_GREY,
                 align=PP_ALIGN.CENTER)
        x += Inches(3.25)

    section_label(slide, "Schema Attributes in Unified Dataset", Inches(2.3))
    attrs_info = [
        ("title",        "MAX of matched titles — prefers longer/more complete form"),
        ("platform",     "Normalized canonical name — entity key component"),
        ("release_date", "MIN of matched dates — earliest known release date"),
        ("developer",    "MAX — longer developer name preferred"),
        ("publisher",    "MAX — longer publisher name preferred"),
        ("genre",        "UNION — genres from all matched sources combined"),
        ("critic_score", "MAX — from DS1 only; missing for DS2/DS3-only records"),
        ("user_score",   "MAX — from DS2 (user_review) or DS3"),
        ("metascore",    "MAX — from DS3 (via Platforms Info 1:n expansion)"),
        ("product_rating","MAX — from DS3 (PEGI/ESRB rating)"),
        ("total_sales",  "MAX — global sales in millions (DS1 only)"),
        ("summary",      "UNION — game description text (DS2 only)"),
        ("source",       "UNION — pipe-separated list of contributing datasets"),
        ("provenance",   "Field-level audit trail: eq/min/max/union/single annotations"),
    ]
    y = Inches(2.65)
    rh = Inches(0.315)
    for i, (attr, desc) in enumerate(attrs_info):
        bg = ROW_A if i % 2 == 0 else ROW_B
        add_rect(slide, Inches(0.35), y, Inches(12.6), rh, bg)
        add_text(slide, attr, Inches(0.5), y + Inches(0.04),
                 Inches(2.0), rh, size=11, bold=True, color=WHITE)
        add_text(slide, desc, Inches(2.6), y + Inches(0.04),
                 Inches(10.1), rh, size=11, color=LIGHT_GREY)
        y += rh


# ── Slide 20 – Summary ───────────────────────────────────────────────────────
def s20_summary(prs):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)
    title_bar(slide, "Summary", "")
    add_rect(slide, 0, Inches(1.05), SLIDE_W, Inches(0.05), ACCENT)

    items = [
        ("Problem",
         "Three heterogeneous Kaggle datasets with overlapping but incomplete video game information",
         ACCENT3),
        ("Entity Key",
         "(title, platform) — same game on different platforms treated as distinct entities",
         ACCENT4),
        ("Schema Mapping",
         "GAV approach: unified target schema + NULL / 1:1 / 1:n mappings per source via apply_mapping()",
         ACCENT),
        ("Normalization",
         "LLM-assisted dicts for platforms/genres/publishers; date parsing; Roman numeral conversion",
         ACCENT),
        ("Blocking",
         "Platform block (level 1) + release-year ±1 block (level 2) — reduces comparisons drastically",
         ACCENT2),
        ("Matching",
         "score = 0.85·lev(title) + 0.15·date_decay — platform mismatch is a hard reject",
         ACCENT2),
        ("Conflict Resolution",
         "Per-field strategies: MIN (dates), MAX (scores/text), UNION (genres/sources), NOTHING (platform)",
         ACCENT3),
        ("Provenance",
         "Every field annotated with eq/min/max/union/single(source) for full transparency",
         ACCENT4),
        ("Result",
         "80 523 unified, sorted, deduplicated records in games_integrated_dataset.csv",
         ACCENT),
    ]

    y = Inches(1.2)
    rh = Inches(0.64)
    for i, (label, desc, color) in enumerate(items):
        bg = RGBColor(0x2C, 0x2C, 0x48) if i % 2 == 0 else RGBColor(0x24, 0x24, 0x3C)
        add_rect(slide, Inches(0.35), y, Inches(12.6), rh, bg)
        add_rect(slide, Inches(0.35), y, Inches(0.1), rh, color)
        add_text(slide, label, Inches(0.55), y + Inches(0.1),
                 Inches(2.2), rh - Inches(0.1), size=12, bold=True, color=color)
        add_text(slide, desc, Inches(2.85), y + Inches(0.1),
                 Inches(9.9), rh - Inches(0.1), size=12, color=LIGHT_GREY)
        y += rh + Inches(0.01)

    add_text(slide,
             "Information and Data Integration  ·  FSU Jena  ·  WS 2025/26",
             Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.3),
             size=10, italic=True, color=SUBTITLE_C, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    s01_title(prs)
    s02_toc(prs)
    s03_motivation(prs)
    s04_datasets(prs)        # all three datasets on one slide
    s06_comparison(prs)      # attribute coverage matrix
    s07_architecture(prs)
    s08_extraction(prs)
    s09_schema_concepts(prs)
    s10_apply_mapping(prs)
    s11_mapping_defs(prs)
    s12_preprocessing(prs)
    s13_llm_mapping(prs)
    s14_matching(prs)
    s15_blocking(prs)
    s16_merging(prs)
    s17_quality_overview(prs)
    s18_quality_issues(prs)
    s19_results(prs)
    s20_summary(prs)

    output = "information_integration_presentation_detailed.pptx"
    prs.save(output)
    print(f"Saved: {output}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
